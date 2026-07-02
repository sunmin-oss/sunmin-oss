#!/usr/bin/env python3
"""
Inspect a PowerPoint file and return structured metadata as JSON.
Usage: python inspect_pptx.py <file.pptx> [--slide N] [--text] [--notes] [--layouts]
"""
import json, sys, argparse
from pathlib import Path


def inspect(filepath, slide_num=None, include_text=False, include_notes=False, include_layouts=False):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu

    prs = Presentation(filepath)
    
    result = {
        "file": str(filepath),
        "slide_width_inches": round(prs.slide_width / 914400, 2),
        "slide_height_inches": round(prs.slide_height / 914400, 2),
        "total_slides": len(prs.slides),
        "slide_layouts_available": [layout.name for layout in prs.slide_layouts],
        "slides": [],
    }
    
    for idx, slide in enumerate(prs.slides):
        if slide_num is not None and idx != slide_num:
            continue
        
        slide_info = {
            "index": idx,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "has_notes": slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip()),
        }
        
        # Shape summary
        shapes = []
        for shape in slide.shapes:
            s = {
                "name": shape.name,
                "type": shape.shape_type.__class__.__name__ if hasattr(shape.shape_type, '__class__') else str(shape.shape_type),
                "left_inches": round(shape.left / 914400, 2) if shape.left else 0,
                "top_inches": round(shape.top / 914400, 2) if shape.top else 0,
                "width_inches": round(shape.width / 914400, 2) if shape.width else 0,
                "height_inches": round(shape.height / 914400, 2) if shape.height else 0,
            }
            
            if shape.has_text_frame:
                s["has_text"] = True
                if include_text:
                    paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        p_info = {"text": para.text}
                        if para.alignment is not None:
                            p_info["alignment"] = str(para.alignment)
                        if para.runs:
                            run = para.runs[0]
                            if run.font.bold:
                                p_info["bold"] = True
                            if run.font.italic:
                                p_info["italic"] = True
                            if run.font.size:
                                p_info["font_size_pt"] = round(run.font.size / 12700, 1)
                            if run.font.name:
                                p_info["font_name"] = run.font.name
                        paragraphs.append(p_info)
                    s["paragraphs"] = paragraphs
            
            if hasattr(shape, "image"):
                try:
                    s["has_image"] = True
                    s["image_content_type"] = shape.image.content_type
                except Exception:
                    pass
            
            if shape.has_table:
                table = shape.table
                s["has_table"] = True
                s["table_rows"] = len(table.rows)
                s["table_cols"] = len(table.columns)
                if include_text:
                    s["table_data"] = [
                        [cell.text for cell in row.cells]
                        for row in table.rows
                    ]
            
            if shape.has_chart:
                s["has_chart"] = True
                s["chart_type"] = str(shape.chart.chart_type)
            
            shapes.append(s)
        
        slide_info["shapes"] = shapes
        
        if include_notes and slide.has_notes_slide:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text
        
        result["slides"].append(slide_info)
    
    if include_layouts:
        layouts = []
        for layout in prs.slide_layouts:
            l = {"name": layout.name, "placeholders": []}
            for ph in layout.placeholders:
                l["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                })
            layouts.append(l)
        result["layout_details"] = layouts
    
    return result


def main():
    parser = argparse.ArgumentParser(description="Inspect PowerPoint file")
    parser.add_argument("file", help="PowerPoint file path")
    parser.add_argument("--slide", type=int, help="Inspect specific slide (0-indexed)")
    parser.add_argument("--text", action="store_true", help="Include text content")
    parser.add_argument("--notes", action="store_true", help="Include speaker notes")
    parser.add_argument("--layouts", action="store_true", help="Include layout details")
    args = parser.parse_args()
    
    if not Path(args.file).exists():
        print(json.dumps({"error": f"File not found: {args.file}"}))
        sys.exit(1)
    
    result = inspect(args.file, args.slide, args.text, args.notes, args.layouts)
    print(json.dumps(result, indent=2, default=str))


if __name__ == "__main__":
    main()
